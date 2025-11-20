# Ref : https://stackoverflow.com/questions/60849601/is-it-possible-to-combine-two-or-more-powerpoints-using-python-pptx
import argparse as argp
import os
import glob
from pptx import Presentation

def main(param=None):
    path = param.path
    print(path)
    pres = glob.glob(f"{path}*.pptx")
    print(pres)
    prs1 = Presentation(pres[0])
    # prs2 = Presentation("Example.pptx")
    for presentation in pres[1:]:
        pres = Presentation(presentation)
        for slide in pres.slides:
            sl = prs1.slides.add_slide(prs1.slide_layouts[1])
            sl.shapes.title.text = slide.shapes.title.text
            try:
                sl.placeholders[1].text = slide.placeholders[1].text
            except:
                sl.placeholders[1].image = slide.placeholders[1].image
        prs1.save(f"{path}merged.pptx")

def env_args():
    # command-line options, argumetns : https://brownbears.tistory.com/413, https://docs.python.org/3/library/argparse.html
    parser = argp.ArgumentParser(description='Merging PowerPoint Presentations in Python')

    parser.add_argument('--path', required=True, default='./', help='Input pptx files(default path=./')

    args = parser.parse_args()
    return args

if __name__ == '__main__':
    args = env_args()
    main(param=args)